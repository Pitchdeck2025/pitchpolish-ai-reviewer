from pptx import Presentation

def extract_pptx_text(file):
    prs = Presentation(file)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        content = {
            "slide_number": i + 1,
            "title": "",
            "body": "",
            "text_count": 0,
            "has_image": False,
            "bullet_count": 0
        }
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not content["title"]:
                    content["title"] = text
                else:
                    content["body"] += text + " "
                    content["text_count"] += len(text.split())
                    content["bullet_count"] += text.count("•")
            if shape.shape_type == 13:
                content["has_image"] = True
        slides_data.append(content)
    return slides_data
